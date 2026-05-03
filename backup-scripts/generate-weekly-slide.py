#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Components Bay — Slide Generator
Design: fond ASM + sidebar navy + tableaux hauteur fixe (noAutofit)
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

import requests, json, os
from datetime import datetime, date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SUPABASE_URL = "https://nwidtkiteamvnomewjux.supabase.co"
SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6Im53aWR0a2l0ZWFtdm5vbWV3anV4Iiwicm9sZSI6ImFub24iLCJpYXQiOjE3NzEzMTM3MDAsImV4cCI6MjA4Njg4OTcwMH0.fvKwq5B8Bdr2Hv67yvB6JRKA2Gu3xTIEgPmcmJj0nvI"

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BG_IMAGE   = os.path.join(SCRIPT_DIR, "asm_slide_bg.png")

OUTPUT_DIR = r"C:\Users\jpellegrini\Desktop\APP 5.5"
if not os.path.exists(OUTPUT_DIR):
    OUTPUT_DIR = r"C:\ComponentsBay_Logs"
if not os.path.exists(OUTPUT_DIR):
    OUTPUT_DIR = SCRIPT_DIR

TODAY  = datetime.now().strftime('%Y-%m-%d')
OUTPUT = os.path.join(OUTPUT_DIR, f"Component_Bay_Slide_{TODAY}.pptx")

# ─── Colors ──────────────────────────────────────────────────────────────────
def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
NAVY    = rgb("1B2B4B"); WHITE   = rgb("FFFFFF"); ROW_ALT = rgb("EEF2F7")
BORDER  = rgb("C8D5E3"); TEXT    = rgb("1E293B"); SUBTEXT = rgb("5A7A99")
AMBER   = rgb("D4900A"); RED_ACC = rgb("C0392B"); BLUE_ACC= rgb("1E6BB5")
CRIMSON = rgb("991B1B"); ORANGE  = rgb("B45309"); GREEN   = rgb("15803D")
MASK_BG = rgb("F2F0EE"); MASK_DK = rgb("1B2A3E")

def i(x): return Inches(x)
def p(x): return Pt(x)

W=i(13.33); H=i(7.5)
SB_X=i(0.32); SB_W=i(2.55); TY=i(1.05)
TTX=SB_X+SB_W+i(0.18); TTW=W-TTX-i(0.32)

# ─── Fetch ────────────────────────────────────────────────────────────────────
def fetch(table):
    hdrs = {"apikey": SUPABASE_KEY, "Authorization": f"Bearer {SUPABASE_KEY}"}
    r = requests.get(f"{SUPABASE_URL}/rest/v1/{table}?select=*", headers=hdrs, timeout=30)
    r.raise_for_status()
    items = []
    for row in r.json():
        d = row.get('data', {})
        if isinstance(d, str): d = json.loads(d)
        items.append(d)
    return items

def days_left(s):
    if not s: return None
    try:
        dt = datetime.strptime(str(s)[:10], '%Y-%m-%d').date()
        return (dt - date.today()).days
    except: return None

def day_color(val):
    try:
        d = int(str(val).replace('d','').strip())
        if d <= 30: return CRIMSON
        if d <= 60: return ORANGE
        return GREEN
    except: return TEXT

def trunc(s, n=32):
    s = str(s or "")
    return s[:n] + "..." if len(s) > n else s

# ─── XML helpers ─────────────────────────────────────────────────────────────
def hex_color(color):
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"

def set_cell_fill(cell, color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', str(color).upper())

def set_cell_border(cell, color, w_pt=0.4):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    w = int(w_pt * 12700)
    for side in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        for old in tcPr.findall(qn(side)): tcPr.remove(old)
        ln = etree.SubElement(tcPr, qn(side))
        ln.set('w', str(w)); ln.set('cap', 'flat'); ln.set('cmpd', 'sng')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr'))
        sc.set('val', str(color).upper())

def set_cell_margins_zero(cell):
    """Remove all internal cell padding so rows can be as small as needed."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # marL, marR, marT, marB = 0
    for attr in ['marL','marR','marT','marB']:
        tcPr.set(attr, '0')

def force_noautofit(cell):
    """Prevent PowerPoint from expanding the row height to fit content."""
    txBody = cell._tc.find(qn('a:txBody'))
    if txBody is None: return
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None: return
    # Remove any existing autofit elements
    for tag in [qn('a:spAutoFit'), qn('a:normAutofit'), qn('a:noAutofit')]:
        for el in bodyPr.findall(tag): bodyPr.remove(el)
    # Add noAutofit — this is the key fix
    etree.SubElement(bodyPr, qn('a:noAutofit'))

def style_cell(cell, text, font_size, bold=False, fg=None, bg=None, align=PP_ALIGN.LEFT):
    if fg is None: fg = TEXT
    # Set text
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = str(text or "")
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = fg
    # Vertical center
    from pptx.enum.text import MSO_ANCHOR
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Background
    if bg: set_cell_fill(cell, bg)
    # Border
    set_cell_border(cell, BORDER)
    # CRITICAL: prevent row auto-expansion
    force_noautofit(cell)
    # CRITICAL: zero cell margins so rows can be compact
    set_cell_margins_zero(cell)

# ─── Shape helpers ────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color: shape.line.color.rgb = line_color; shape.line.width = p(0.5)
    else: shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size, bold=False, color=None, align=PP_ALIGN.LEFT):
    if color is None: color = WHITE
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    para = tf.paragraphs[0]; para.alignment = align
    run = para.add_run(); run.text = str(text)
    run.font.size = size; run.font.bold = bold
    run.font.color.rgb = color; run.font.name = "Calibri"
    return tb

# ─── Slide builder ────────────────────────────────────────────────────────────
def build_slide(prs, accent, title, subtitle, header_label,
                count, headers, rows, col_w_in, days_col=-1):

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    if os.path.exists(BG_IMAGE):
        slide.shapes.add_picture(BG_IMAGE, i(0), i(0), W, H)

    # Mask template text
    add_rect(slide, i(0),    i(0.60), i(6.5),  i(0.55), MASK_BG)
    add_rect(slide, i(12.5), i(7.10), i(0.85), i(0.40), MASK_DK)
    add_rect(slide, TTX,     TY,      TTW,      H-TY-i(0.20), MASK_BG)

    # Sidebar
    sb_h = H - TY - i(0.68)
    add_rect(slide, SB_X, TY, SB_W, sb_h, NAVY)
    add_rect(slide, SB_X, TY, SB_W, i(0.07), accent)
    add_text(slide, str(count), SB_X, TY+i(0.12), SB_W, i(1.35), p(72), bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_rect(slide, SB_X+i(0.2), TY+i(1.55), SB_W-i(0.4), i(0.02), accent)
    add_text(slide, title.upper(), SB_X, TY+i(1.62), SB_W, i(1.15), p(17), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    PY = TY+i(2.95)
    add_rect(slide, SB_X+i(0.2), PY, SB_W-i(0.4), i(0.52), accent)
    add_text(slide, subtitle, SB_X+i(0.2), PY, SB_W-i(0.4), i(0.52), p(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, TODAY, SB_X, H-i(0.88), SB_W, i(0.22), p(7.5), color=SUBTEXT, align=PP_ALIGN.CENTER)

    # Header bar
    BH = i(0.40)
    add_rect(slide, TTX, TY, TTW, BH, accent)
    add_text(slide, f"!  {header_label}", TTX+i(0.12), TY, TTW-i(1.6), BH, p(9.5), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(slide, TTX+TTW-i(1.1), TY+i(0.06), i(0.95), BH-i(0.12), WHITE, WHITE)
    add_text(slide, str(count), TTX+TTW-i(1.1), TY+i(0.06), i(0.95), BH-i(0.12), p(13), bold=True, color=accent, align=PP_ALIGN.CENTER)

    # Table dimensions — STRICT: never exceed slide
    HDR_H   = i(0.26)
    TBL_Y   = TY + BH
    MAX_H   = H - TBL_Y - i(0.75)   # hard limit
    n_rows  = max(len(rows), 1)
    n_cols  = len(headers)
    cw_emu  = [i(w) for w in col_w_in]

    # Row height: divide available space exactly
    avail   = MAX_H - HDR_H
    data_rh = int(avail / n_rows)
    # Cap for sparse tables
    data_rh = min(data_rh, i(0.26))
    TBL_H   = HDR_H + data_rh * n_rows

    # Font scales with row height
    row_in  = data_rh / 914400
    font_sz = p(6.5) if row_in < 0.155 else (p(7.0) if row_in < 0.195 else p(7.5))

    # Create table
    tbl_shape = slide.shapes.add_table(n_rows+1, n_cols, TTX, TBL_Y, TTW, TBL_H)
    tbl = tbl_shape.table

    # Set column widths
    for ci, cw in enumerate(cw_emu):
        tbl.columns[ci].width = cw

    # Set row heights
    tbl.rows[0].height = HDR_H
    for ri in range(1, n_rows+1):
        tbl.rows[ri].height = data_rh

    # Header row
    for ci, hdr in enumerate(headers):
        style_cell(tbl.cell(0, ci), hdr, p(8), bold=True, fg=AMBER, bg=NAVY)

    # Data rows
    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else ROW_ALT
        for ci, val in enumerate(row):
            txt = str(val or "").strip() or " "
            fg = TEXT; bold = False
            if ci == days_col and txt:
                fg = day_color(txt); bold = True
            style_cell(tbl.cell(ri+1, ci), txt, font_sz, bold=bold, fg=fg, bg=bg)

    # Bottom accent line — fixed at bottom of table area
    add_rect(slide, TTX, TBL_Y+TBL_H-i(0.03), TTW, i(0.025), accent)

# ─── Main ─────────────────────────────────────────────────────────────────────
def main():
    print("="*55)
    print("  Components Bay — Slide Generator")
    print(f"  {TODAY}")
    print("="*55)
    print("\nChargement Supabase...")
    try:
        efs=fetch('efs'); efs_cyl=fetch('efs_cylinders'); liferafts=fetch('liferafts')
        maint=fetch('maintenance'); comp=fetch('composite'); avio=fetch('avionic')
        eng=fetch('engine'); rotor=fetch('rotorbay'); iaft=fetch('iafteaft')
        pol=fetch('pol'); tools=fetch('tools'); troop=fetch('troopseats')
    except Exception as e:
        print(f"ERREUR Supabase: {e}"); raise

    s1=[]
    for mod, items in [("Maintenance",maint),("Composite",comp),("Avionic",avio),
                        ("Engine",eng),("Rotor Bay",rotor),("IAFT/EAFT",iaft),
                        ("POL",pol),("Tools",tools),("Troop Seat",troop),
                        ("EFS Float",efs),("EFS Cyl.",efs_cyl)]:
        for it in items:
            if it.get('serviceability')=='Unserviceable':
                pn=it.get('pnWheel') or it.get('partNumber','')
                s1.append([mod,trunc(it.get('designation',''),28),trunc(pn,18),
                           it.get('serialNumber',''),trunc(it.get('reason','Unserviceable'),28)])

    s2=[]
    for it in liferafts:
        if it.get('serviceability')=='Unserviceable':
            s2.append([it.get('partNumber',''),it.get('serialNumber',''),
                       trunc(it.get('reason','Unserviceable'),36)])

    s3=[]
    for it in efs:
        if it.get('serviceability')=='Unserviceable': continue
        best=999
        for f in ['next18M','next36M']:
            d=days_left(it.get(f))
            if d is not None and 0<=d<=90 and d<best: best=d
        if best<999:
            s3.append([it.get('hc',''),trunc(it.get('designation',''),22),
                       it.get('partNumber',''),it.get('serialNumber',''),
                       it.get('next18M',''),it.get('next36M',''),f"{best}d"])
    s3.sort(key=lambda r:int(r[6].replace('d','')))

    s4=[]
    for it in efs_cyl:
        if it.get('serviceability')=='Unserviceable': continue
        best=999
        for f in ['next18M','next60M']:
            d=days_left(it.get(f))
            if d is not None and 0<=d<=90 and d<best: best=d
        if best<999:
            s4.append([it.get('hc',''),trunc(it.get('designation',''),22),
                       it.get('partNumber',''),it.get('serialNumber',''),
                       it.get('next18M',''),it.get('next60M',''),f"{best}d"])
    s4.sort(key=lambda r:int(r[6].replace('d','')))

    print(f"  Slide 1 — Other Parts U/S:        {len(s1)}")
    print(f"  Slide 2 — Life Raft U/S:           {len(s2)}")
    print(f"  Slide 3 — EFS <90d:                {len(s3)}")
    print(f"  Slide 4 — EFS Cylinders <90d:      {len(s4)}")
    print("\nGeneration PPTX...")

    prs = Presentation()
    prs.slide_width=W; prs.slide_height=H

    cw1=[1.30,2.20,1.95,0.90,3.61]
    cw2=[1.85,3.22,4.89]
    cw34=[1.30,2.05,1.90,0.78,1.35,1.35,1.23]

    build_slide(prs,RED_ACC,"Other Parts","NEED TO BE C/OUT","ALL UNSERVICEABLE ITEMS (ORDERED)",len(s1),["MODULE","DESIGNATION","P/N","S/N","REASON"],s1,cw1)
    build_slide(prs,RED_ACC,"Life Raft","NEED TO BE C/OUT","LIFE RAFT — NEED TO BE C/OUT",len(s2),["P/N","S/N","REASON"],s2,cw2)
    build_slide(prs,BLUE_ACC,"EFS","DUE WITHIN 90 DAYS","SERVICEABLE EFS — DUE WITHIN 90 DAYS",len(s3),["H/C","DESIGNATION","P/N","S/N","NEXT 18M","NEXT 36M","DAYS LEFT"],s3,cw34,days_col=6)
    build_slide(prs,BLUE_ACC,"EFS Cylinders","DUE WITHIN 90 DAYS","SERVICEABLE EFS CYLINDERS — DUE WITHIN 90 DAYS",len(s4),["H/C","DESIGNATION","P/N","S/N","NEXT 18M","NEXT 60M","DAYS LEFT"],s4,cw34,days_col=6)

    prs.save(OUTPUT)
    print(f"\n✅  {OUTPUT}")
    print("="*55)

if __name__ == '__main__':
    main()
